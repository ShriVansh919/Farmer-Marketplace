#!/usr/bin/env python3
"""Generate Kisaan Bazaar project docs (proposal, final report, presentation)."""
import os

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB

OUT = os.path.dirname(os.path.abspath(__file__))
GREEN = RGBColor(0x2E, 0x7D, 0x32)


def heading(doc, text, size=14):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(size)
    r.font.color.rgb = GREEN


def body(doc, text, bold=False, align=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = bold
    if align == "center":
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return p


def bullets(doc, items):
    for it in items:
        p = doc.add_paragraph(it, style="List Bullet")
        p.paragraph_format.space_after = Pt(2)


def table(doc, rows, header=None):
    t = doc.add_table(rows=1, cols=len(rows[0]))
    t.style = "Light Shading Accent 1"
    for i, h in enumerate((header or rows[0])):
        t.rows[0].cells[i].text = h
    data = rows if header else rows[1:]
    for r in data:
        cells = t.add_row().cells
        for i, v in enumerate(r):
            cells[i].text = str(v)
    return t


TEAM = [
    ("Person 1", "Authentication", "Splash, Login, Register, Logout, SharedPreferences"),
    ("Person 2", "Home", "Dashboard, Search Bar, Categories, Featured Products, Bottom Navigation"),
    ("Person 3", "Products", "Product List, Product Details, Dummy JSON, Product Model"),
    ("Person 4", "Cart & Orders", "Cart, Checkout, Orders, Quantity +/-, Total (Provider)"),
    ("Person 5", "Profile & Integration", "Profile, Settings, Theme, Final Testing, GitHub merge"),
]


def build_proposal():
    doc = Document()
    body(doc, "PROJECT PROPOSAL", bold=True)
    body(doc, "KISAN BAZAAR — A Flutter-Based Digital Marketplace for Agricultural Produce", bold=True)
    body(doc, "Classroom Project — Flutter Hybrid Mobile Application Development (JOVAC)", align="center")
    doc.add_paragraph()

    heading(doc, "1. Introduction")
    body(doc, "Agricultural markets involve multiple intermediaries between farmers and consumers, "
              "reducing transparency and the share of the final price received by farmers. Kisan "
              "Bazaar is a mobile-based agricultural marketplace built with Flutter where products "
              "such as vegetables, fruits, grains, and dairy can be displayed, searched, "
              "categorised, added to a cart, and ordered.")

    heading(doc, "2. Problem Statement")
    bullets(doc, [
        "Dependence on intermediaries.",
        "Limited visibility of available agricultural products.",
        "Difficulty in comparing products and prices.",
        "No lightweight digital platform tailored to agricultural produce.",
    ])

    heading(doc, "3. Proposed Solution")
    body(doc, "A Flutter app providing: user authentication, product catalogue, categories, search, "
              "product details, shopping cart, quantity management, order management, and profile "
              "& settings. Products are loaded from structured JSON; session state is persisted "
              "with SharedPreferences; cart state is shared with Provider.")

    heading(doc, "4. Objectives")
    bullets(doc, [
        "Develop an intuitive Flutter-based mobile interface.",
        "Provide secure user registration and login.",
        "Organise agricultural products into categories with search.",
        "Support cart with quantity controls and dynamic total.",
        "Maintain order history and a profile/settings section.",
        "Develop collaboratively using Git/GitHub.",
    ])

    heading(doc, "5. Technology Stack")
    doc.add_paragraph("Flutter, Dart, Material Design, Provider, SharedPreferences, JSON, Git/GitHub.",
                      style="List Bullet")

    heading(doc, "6. Team Structure (5 Members)")
    table(doc, TEAM, header=["Member", "Responsibility", "Main Screens"])

    heading(doc, "7. Future Scope")
    body(doc, "Farmer dashboard, real-time inventory, online payments, delivery tracking, "
              "reviews and ratings, multilingual support, backend/cloud deployment.")
    doc.save(os.path.join(OUT, "Project_Proposal.docx"))


def build_report():
    doc = Document()
    body(doc, "FINAL PROJECT REPORT", bold=True)
    body(doc, "KISAN BAZAAR — A Flutter-Based Digital Agricultural Marketplace", bold=True)
    body(doc, "Classroom Project — Flutter Hybrid Mobile Application Development", align="center")
    doc.add_paragraph()

    heading(doc, "1. Abstract")
    body(doc, "Kisan Bazaar is a Flutter-based mobile marketplace for agricultural products. The "
              "application provides authentication, product browsing via categorised listings, "
              "keyword search, a shopping cart with quantity and price management, order history, "
              "and a user profile — all implemented with a modular, widget-based architecture. "
              "Cart state is shared through the Provider package, session data is persisted with "
              "SharedPreferences, and product information is parsed from structured JSON. Git and "
              "GitHub were used for version control and collaborative development.")

    heading(doc, "2. Architecture")
    bullets(doc, [
        "Presentation layer: Flutter widgets/screens.",
        "State management: Provider (shared cart and session state).",
        "Data layer: local JSON product data + SharedPreferences persistence.",
        "Workflow: Splash → Login/Register → Home → Products → Cart → Checkout → Orders → Profile.",
    ])

    heading(doc, "3. Modules and Screens")
    table(doc, TEAM, header=["Member", "Responsibility", "Main Screens"])
    doc.add_paragraph()
    bullets(doc, [
        "Splash: brand screen, restores previous session.",
        "Login/Register: validated forms; invalid credentials show an error.",
        "Home: branding header, search bar, promotional banner, horizontal categories, featured grid, bottom navigation.",
        "Search: real-time keyword filtering of the catalogue.",
        "Categories: All, Vegetables, Fruits, Grains, Dairy, Organic.",
        "Product Details: description, price, rating, stock, farmer and location.",
        "Cart: add/remove, quantity +/-, subtotal/total recalculation (Provider).",
        "Orders: order history with status after checkout.",
        "Profile: user details, dark-theme toggle, logout.",
    ])

    heading(doc, "4. Testing (manual, verified on running app)")
    table(doc, [
        ("Login", "Valid credentials", "Pass"),
        ("Login", "Invalid credentials", "Pass"),
        ("Register", "New account + duplicate email", "Pass"),
        ("Validation", "Empty/invalid fields", "Pass"),
        ("Search", "\"apple\"", "Pass"),
        ("Category", "Vegetables filter", "Pass"),
        ("Add Cart", "Select product", "Pass"),
        ("Quantity", "+ / - buttons", "Pass"),
        ("Total", "Cart changes", "Pass"),
        ("Checkout", "Place order", "Pass"),
        ("Orders", "Order history", "Pass"),
        ("Theme", "Dark mode toggle", "Pass"),
        ("Logout", "Session cleared", "Pass"),
    ], header=["Test Case", "Input / Action", "Status"])

    heading(doc, "5. Future Scope")
    body(doc, "Farmer dashboard, backend REST API, online payments, delivery tracking, "
              "reviews and ratings, multilingual support, Google Maps integration.")
    doc.save(os.path.join(OUT, "Final_Report.docx"))


def build_ppt():
    prs = Presentation()
    prs.slide_width = PIn(13.333)
    prs.slide_height = PIn(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title, lines):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(PIn(0.8), PIn(0.5), PIn(11.7), PIn(1.2))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.size = PPt(40)
        r.font.bold = True
        r.font.color.rgb = PRGB(0x2E, 0x7D, 0x32)
        b = slide.shapes.add_textbox(PIn(1.0), PIn(2.0), PIn(11.3), PIn(4.8))
        first = True
        for line in lines:
            p = b.text_frame.paragraphs[0] if first else b.text_frame.add_paragraph()
            first = False
            r = p.add_run()
            r.text = line
            r.font.size = PPt(20)
            p.space_after = PPt(10)

    add_slide("Kisan Bazaar", [
        "A Flutter marketplace connecting farmers and buyers",
        "Browse → Search → Cart → Checkout → Orders",
        "Flutter + Provider + SharedPreferences + JSON",
    ])
    add_slide("Problem", [
        "Multiple middlemen cut farmer margins",
        "Buyers struggle to discover local produce",
        "Low price transparency, fragmented information",
    ])
    add_slide("Solution", [
        "Mobile-first marketplace for vegetables, fruits, grains, dairy",
        "Categorised catalogue with keyword search",
        "Cart with live price calculation and order history",
    ])
    add_slide("User Workflow", [
        "Splash → Login/Register → Home",
        "Search / Categories → Product List → Product Details",
        "Add to Cart → Checkout → Orders → Profile",
    ])
    add_slide("Team Structure", [
        "1 · Authentication — Splash, Login, Register, Logout",
        "2 · Home — Dashboard, Search, Categories, Bottom Nav",
        "3 · Products — List, Details, JSON, Product Model",
        "4 · Cart & Orders — Cart, Checkout, Orders (Provider)",
        "5 · Profile & Integration — Profile, Theme, Git, Testing",
    ])
    add_slide("Testing", [
        "Manual matrix — 13/13 test cases Pass",
        "flutter analyze — no issues",
        "Widget + unit tests — all passing",
    ])
    add_slide("Future Scope", [
        "Farmer dashboard & real-time inventory",
        "REST API backend + payments",
        "Reviews, multilingual, Google Maps",
    ])
    add_slide("Thank You", ["Questions welcome", "Kisan Bazaar Team — GLA University, Mathura"])
    prs.save(os.path.join(OUT, "Presentation.pptx"))


# More deliverable docs that don't need the docx/pptx writers: demo script and viva prep.
DEMO_SCRIPT = """# Kisan Bazaar — Demo Video Script

Record your screen (phone emulator or real device) while this narration plays.
Each shot under 15 seconds. Total ≈ 2 minutes.

**Shot 1 — Intro (5s)** · Splash screen.
*"This is Kisan Bazaar, a Flutter marketplace for agricultural produce."*

**Shot 2 — Register (15s)** · Register form: name, email, password.
*"New users register with a validated form — empty or invalid fields are blocked."*

**Shot 3 — Login + session (10s)** · Log in with the demo account, reopen the app.
*"Login is persisted with SharedPreferences, so the session survives app restarts."*

**Shot 4 — Home (12s)** · Brand header, banner, categories, featured grid.
*"The home screen shows a promo banner, category chips, and the featured product grid."*

**Shot 5 — Search (8s)** · Type "apple".
*"The search tab filters products live as you type."*

**Shot 6 — Product details (12s)** · Open a product, add to cart.
*"Product details show price, rating, stock, farmer and location."*

**Shot 7 — Cart (15s)** · Change quantity +/-, watch total update, remove item.
*"The cart updates quantities and totals instantly via Provider state."*

**Shot 8 — Checkout + Orders (12s)** · Place order, open Orders tab.
*"Checkout creates an order that appears in the Orders history."*

**Shot 9 — Profile + theme (8s)** · Toggle dark theme, log out.
*"The profile shows user details, a theme toggle, and logout."*

**Shot 10 — Tests (10s)** · Terminal: `flutter test` + `flutter analyze`.
*"All widget tests pass and static analysis is clean."*
"""

VIVA_PREP = """# Kisan Bazaar — Viva Preparation

## Flutter / Dart
1. **What is Flutter?** — Google's open-source UI toolkit; one Dart codebase compiles to native
   Android/iOS and desktop.
2. **Stateless vs Stateful?** — Stateless is immutable UI; Stateful holds mutable state
   (controllers, selected category, etc.).
3. **What is hot reload?** — Applies code changes live while preserving state.

## State management (Provider)
4. **Why Provider?** — Shared, observable state. `ChangeNotifier` + `Consumer`/`watch` rebuild
   only the widgets that listen.
5. **Which providers?** — AuthProvider (session), CartProvider (cart + totals),
   OrdersProvider (order history), ThemeProvider (dark mode).
6. **`context.read` vs `watch`?** — `watch` rebuilds on change; `read` fetches once.

## Forms & validation
7. **How is the form validated?** — `Form` + `GlobalKey<FormState>`; each `TextFormField` has a
   `validator`; `validate()` gates submission.
8. **Validations used?** — required fields, email regex, min password length, password match.

## Navigation
9. **How do you navigate?** — `Navigator.push` (product detail), `pushAndRemoveUntil` (login→home),
   and bottom `NavigationBar` tabs with `IndexedStack`.

## Local storage
10. **What is SharedPreferences?** — Key-value local store; keeps the logged-in user and order
    history across sessions. Demo account is seeded on first run.

## Products
11. **Where do products come from?** — `assets/json/products.json` (26 products), read with
    `rootBundle` and parsed into the `Product` model via `FutureBuilder`.
12. **Product fields?** — id, name, price, category, description, rating, stock, image, imageUrl,
    farmerName, location.

## Cart & Orders
13. **How does the cart work?** — CartProvider holds `CartItem`s; add/increment (capped at stock),
    decrement/remove; subtotal/total recomputed on every change.
14. **Checkout?** — Confirmation dialog → `OrdersProvider.placeOrder()` → order list persisted.

## Error handling
15. **How do you handle failures?** — Loading indicators, friendly error view with Retry for
    product loading, validation errors and SnackBars for auth issues.

## Git / team work
16. **How is the team organised?** — Five modules with clear owners (Auth, Home, Products,
    Cart & Orders, Profile & Integration); feature branches merged to `main` on GitHub.
17. **Commit style?** — Conventional commits: `feat(cart): add quantity controls`.

## Project-specific
18. **Why a marketplace?** — Reduce middlemen, improve price transparency, and give a simple
    mobile-first way to buy farm produce.
19. **Number of screens?** — Splash, Login, Register, Home, Search, Product Details, Cart,
    Orders, Profile (9 total).
20. **What would you add?** — Farmer dashboard, REST backend, payments, delivery tracking,
    multilingual support.
"""


if __name__ == "__main__":
    build_proposal()
    build_report()
    build_ppt()
    with open(os.path.join(OUT, "DEMO_SCRIPT.md"), "w") as f:
        f.write(DEMO_SCRIPT)
    with open(os.path.join(OUT, "VIVA_PREP.md"), "w") as f:
        f.write(VIVA_PREP)
    print("Docs generated in", OUT)